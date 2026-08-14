import io
import os
from pypdf import PdfReader
from pptx import Presentation
from docx import Document

def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """
    Trích xuất text sạch từ các file thô: PDF, PPTX, DOCX, TXT.
    """
    ext = os.path.splitext(filename)[1].lower()
    extracted_text = ""

    try:
        if ext == ".pdf":
            reader = PdfReader(io.BytesIO(file_bytes))
            pages_text = []
            for i, page in enumerate(reader.pages):
                txt = page.extract_text()
                if txt:
                    pages_text.append(f"--- Trang {i+1} ---\n{txt}")
            extracted_text = "\n".join(pages_text)

        elif ext == ".pptx":
            prs = Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_words = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_words.append(shape.text)
                if slide_words:
                    slides_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_words))
            extracted_text = "\n".join(slides_text)

        elif ext in [".docx", ".doc"]:
            doc = Document(io.BytesIO(file_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            extracted_text = "\n".join(paragraphs)

        elif ext == ".txt":
            extracted_text = file_bytes.decode("utf-8", errors="ignore")

        else:
            extracted_text = f"Định dạng {ext} chưa được hỗ trợ."

    except Exception as e:
        extracted_text = f"Lỗi đọc file {filename}: {str(e)}"

    # Làm sạch văn bản thừa
    cleaned_lines = [line.strip() for line in extracted_text.splitlines() if line.strip()]
    return "\n".join(cleaned_lines)
